import io
import pathlib
import logging
from typing import Optional

logger = logging.getLogger(__name__)

def extract_text_from_drive(service, file_id, mime_type) -> Optional[str]:
    """Download a Drive file and return its plain-text content, or None on failure."""
    try:
        # 1. Native Google Docs -> Export to Text
        if mime_type == "application/vnd.google-apps.document":
            raw = service.files().export(fileId=file_id, mimeType="text/plain").execute()
            return raw.decode("utf-8", errors="ignore") if isinstance(raw, bytes) else str(raw)

        # 2. Native Google Sheets -> Export to XLSX -> Parse
        if mime_type == "application/vnd.google-apps.spreadsheet":
            raw = service.files().export(fileId=file_id, mimeType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet").execute()
            return _extract_spreadsheet_from_bytes(raw)

        # 3. Native Google Slides -> Export to PPTX -> Parse
        if mime_type == "application/vnd.google-apps.presentation":
            raw = service.files().export(fileId=file_id, mimeType="application/vnd.openxmlformats-officedocument.presentationml.presentation").execute()
            return _extract_pptx_from_bytes_raw(raw)

        # 4. Standard Text/CSV
        if mime_type in ("text/plain", "text/csv"):
            from googleapiclient.http import MediaIoBaseDownload
            buf = io.BytesIO()
            dl = MediaIoBaseDownload(buf, service.files().get_media(fileId=file_id))
            done = False
            while not done:
                _, done = dl.next_chunk()
            return buf.getvalue().decode("utf-8", errors="ignore")

        # 5. Office Formats
        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return _extract_docx_from_bytes(service, file_id)

        if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or mime_type == "application/vnd.ms-excel":
            return _extract_spreadsheet_from_bytes(service=service, file_id=file_id)

        if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return _extract_pptx_from_bytes(service, file_id)

        # 6. PDF
        if mime_type == "application/pdf":
            return _extract_pdf_from_bytes(service, file_id)

    except Exception as exc:
        logger.warning(f"extract_text_from_drive failed for {file_id}: {exc}")
    return None

def extract_text_from_local(filepath: pathlib.Path) -> Optional[str]:
    """Robust extraction from a local file, ensuring zero garbage text fallbacks."""
    ext = filepath.suffix.lower()
    try:
        if ext in (".txt", ".md", ".csv"):
            return filepath.read_text(encoding="utf-8", errors="ignore").strip()

        if ext == ".docx":
            return _extract_docx_local(filepath)

        if ext == ".pdf":
            return _extract_pdf_local(filepath)

        if ext in (".xlsx", ".xls"):
            return _extract_spreadsheet_local(filepath)

        if ext == ".pptx":
            return _extract_pptx_local(filepath)

    except Exception as exc:
        logger.error(f"Robust local extraction failed for {filepath.name}: {exc}")
    return None

# --- Specialized Internal Parsers ---

def _extract_pdf_local(filepath: pathlib.Path) -> Optional[str]:
    """Highly robust PDF extraction using PyMuPDF (fitz) with PyPDF fallback."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(str(filepath))
        text = "\n".join([page.get_text() for page in doc])
        doc.close()
        if text.strip():
            return text.strip()
    except Exception as e:
        logger.debug(f"PyMuPDF failed for {filepath}: {e}")

    try:
        from pypdf import PdfReader
        reader = PdfReader(str(filepath))
        text = "\n".join([page.extract_text() or "" for page in reader.pages])
        if text.strip():
            return text.strip()
    except Exception as e:
        logger.debug(f"PyPDF fallback failed for {filepath}: {e}")

    return None

def _extract_docx_local(filepath: pathlib.Path) -> Optional[str]:
    """Structure-preserving extraction from .docx using python-docx."""
    try:
        import docx as docx_lib
        doc = docx_lib.Document(str(filepath))
        return _docx_to_structured_text(doc)
    except Exception as e:
        logger.warning(f"DOCX extraction failed: {e}")
        return None

def _extract_pptx_local(filepath: pathlib.Path) -> Optional[str]:
    """Extracts text from all slides in a .pptx file."""
    try:
        from pptx import Presentation
        prs = Presentation(str(filepath))
        parts = []
        for i, slide in enumerate(prs.slides):
            parts.append(f"\n--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n".join(parts).strip() if parts else None
    except Exception as e:
        logger.warning(f"PPTX extraction failed: {e}")
        return None

def _docx_to_structured_text(doc) -> Optional[str]:
    """Shared helper: converts a python-docx Document into structured text."""
    parts = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        style_name = (p.style.name or "").lower()
        if style_name.startswith("heading"):
            try:
                level = int(style_name.replace("heading", "").strip())
                level = min(max(level, 1), 6)
            except ValueError:
                level = 2
            parts.append(f"\n{'#' * level} {text}")
        elif style_name == "title":
            parts.append(f"\n# {text}")
        else:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                parts.append(" | ".join(row_text))
    return "\n".join(parts).strip() if parts else None

def _extract_spreadsheet_local(filepath: pathlib.Path) -> Optional[str]:
    """Extract rows from spreadsheets into a structured string format."""
    try:
        import pandas as pd
        df_dict = pd.read_excel(str(filepath), sheet_name=None)
        return _pandas_dict_to_text(df_dict)
    except Exception as e:
        logger.warning(f"Spreadsheet extraction failed: {e}")
        return None

def _pandas_dict_to_text(df_dict) -> str:
    all_text = []
    for sheet_name, df in df_dict.items():
        all_text.append(f"Sheet: {sheet_name}")
        all_text.append(df.to_string(index=False))
    return "\n\n".join(all_text).strip()

# --- Byte-based Drive helpers ---

def _extract_pdf_from_bytes(service, file_id) -> Optional[str]:
    from googleapiclient.http import MediaIoBaseDownload
    buf = io.BytesIO()
    dl = MediaIoBaseDownload(buf, service.files().get_media(fileId=file_id))
    done = False
    while not done:
        _, done = dl.next_chunk()
    buf.seek(0)
    try:
        import fitz
        doc = fitz.open(stream=buf.read(), filetype="pdf")
        text = "\n".join([page.get_text() for page in doc])
        doc.close()
        return text.strip() if text.strip() else None
    except Exception:
        return None

def _extract_docx_from_bytes(service, file_id) -> Optional[str]:
    from googleapiclient.http import MediaIoBaseDownload
    buf = io.BytesIO()
    dl = MediaIoBaseDownload(buf, service.files().get_media(fileId=file_id))
    done = False
    while not done:
        _, done = dl.next_chunk()
    buf.seek(0)
    try:
        import docx as docx_lib
        doc = docx_lib.Document(buf)
        return _docx_to_structured_text(doc)
    except Exception:
        return None

def _extract_spreadsheet_from_bytes(raw_bytes=None, service=None, file_id=None) -> Optional[str]:
    try:
        if raw_bytes:
            buf = io.BytesIO(raw_bytes)
        else:
            from googleapiclient.http import MediaIoBaseDownload
            buf = io.BytesIO()
            dl = MediaIoBaseDownload(buf, service.files().get_media(fileId=file_id))
            done = False
            while not done:
                _, done = dl.next_chunk()
        buf.seek(0)
        import pandas as pd
        df_dict = pd.read_excel(buf, sheet_name=None)
        return _pandas_dict_to_text(df_dict)
    except Exception:
        return None

def _extract_pptx_from_bytes(service, file_id) -> Optional[str]:
    from googleapiclient.http import MediaIoBaseDownload
    buf = io.BytesIO()
    dl = MediaIoBaseDownload(buf, service.files().get_media(fileId=file_id))
    done = False
    while not done:
        _, done = dl.next_chunk()
    return _extract_pptx_from_bytes_raw(buf.getvalue())

def _extract_pptx_from_bytes_raw(raw_bytes) -> Optional[str]:
    try:
        from pptx import Presentation
        buf = io.BytesIO(raw_bytes)
        prs = Presentation(buf)
        parts = []
        for i, slide in enumerate(prs.slides):
            parts.append(f"\n--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n".join(parts).strip() if parts else None
    except Exception:
        return None
